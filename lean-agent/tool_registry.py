"""
tool_registry.py – All agent tools (v2.0)
=========================================
- Filesystem: read, write, edit (with fuzzy matching), copy, move, delete, tree
- Shell: sync/async command execution, background tasks, live streaming
- Web: search, read URL, browser automation
- Knowledge: learn/recall patterns, key‑value memory, git
- Type‑hinted, robust error handling, fully tested
"""

from __future__ import annotations
import asyncio
import subprocess
import os
import json
import re
import shutil
import base64
import threading
import uuid
import datetime
import sqlite3
import difflib
from typing import Any, Callable, Dict, List, Optional, Tuple
from urllib.parse import quote_plus as url_quote

import httpx
import platform
from pathlib import Path
import urllib.request
import urllib.parse
from urllib.error import URLError, HTTPError
import zipfile
import tarfile
import time
import hashlib

# Try to import optional dependencies for extended tools
HAVE_PSUTIL = False
HAVE_GIT = False
HAVE_PANDAS = False
HAVE_MATPLOTLIB = False
HAVE_JEDI = False

try:
    import psutil
    HAVE_PSUTIL = True
except ImportError:
    pass

try:
    import git
    HAVE_GIT = True
except ImportError:
    pass

try:
    import pandas as pd
    HAVE_PANDAS = True
except ImportError:
    pass

try:
    import matplotlib.pyplot as plt
    HAVE_MATPLOTLIB = True
except ImportError:
    pass

try:
    import jedi
    HAVE_JEDI = True
except ImportError:
    pass

# ─── Verification & Output Formatting Helpers ─────────────────────────────────

def truncate_stdout(stdout_str: str, max_lines: int = 300) -> str:
    """Truncate stdout to first N and last N lines if it exceeds max_lines."""
    if not stdout_str:
        return ""
    lines = stdout_str.splitlines(keepends=True)
    if len(lines) <= max_lines:
        return stdout_str
    half = max_lines // 2
    first_part = "".join(lines[:half])
    last_part = "".join(lines[-half:])
    excess = len(lines) - max_lines
    return f"{first_part}\n... [TRUNCATED {excess} LINES OF OUTPUT FOR CONTEXT SANITY] ...\n\n{last_part}"

def check_python_syntax(content: str) -> Optional[str]:
    """Built-in AST compilation check for Python scripts."""
    try:
        compile(content, "<string>", "exec")
        return None
    except SyntaxError as e:
        return f"Python SyntaxError on line {e.lineno}: {e.msg}\nCode snippet:\n{e.text}"
    except Exception as e:
        return f"Syntax compilation check failed: {e}"

def check_json_syntax(content: str) -> Optional[str]:
    """JSON validation check."""
    try:
        json.loads(content)
        return None
    except json.JSONDecodeError as e:
        return f"JSON DecodeError: {e.msg} at line {e.lineno}, column {e.colno}"

# ---------------------------------------------------------------------------
# Globals & workspace
# ---------------------------------------------------------------------------

def _get_initial_workspace() -> str:
    """Get initial workspace WITHOUT resolving junctions."""
    # Don't use abspath - it resolves junctions!
    script_dir = os.path.dirname(__file__)
    parent_dir = os.path.dirname(script_dir)
    
    # Use normpath only (doesn't resolve junctions)
    candidate = os.path.normpath(parent_dir)
    
    # Check if this looks like a physical path that should be user-facing
    if re.search(r'(?i)[/\\\\]s([/\\\\]|$)', candidate):
        user_path = os.path.normpath(re.sub(r'(?i)([/\\\\])s([/\\\\]|$)', r'\1Downloads\2', candidate))
        if os.path.exists(user_path):
            return user_path
    
    return candidate


def _decode_base64(data: str) -> str:
    """Decode base64 string to original text, preserving all whitespace."""
    try:
        return base64.b64decode(data).decode('utf-8')
    except Exception as e:
        return f"[Error decoding base64: {e}]"

# ─── Path Mapping Globals ──────────────────────────────────────────────────
CURRENT_WORKSPACE = os.getcwd()
LOGICAL_ROOT = None
PHYSICAL_ROOT = None

_file_cache: Dict[str, Tuple[float, str]] = {}   # path → (mtime, text)

IGNORED_TREE_DIRS = set()
TRIMMED_TREE_DIRS = set()
MAX_TREE_NODES = 10000

SAFE_COMMANDS = {
    "dir", "ls", "pwd", "echo", "git status", "git log",
    "python --version", "pip --version", "npm --version", "node --version"
}

# Background process tracking
active_processes: Dict[str, Dict[str, Any]] = {}

# Lazy import of browser manager
browser_manager = None
def _get_browser_manager():
    global browser_manager
    if browser_manager is None:
        from browser_manager import browser_manager as bm
        browser_manager = bm
    return browser_manager


# ---------------------------------------------------------------------------
# Workspace / DB helpers
# ---------------------------------------------------------------------------
def _init_db():
    conn = sqlite3.connect(os.path.join(CURRENT_WORKSPACE, "memories.db"))
    conn.execute("""CREATE TABLE IF NOT EXISTS memories (
        id INTEGER PRIMARY KEY,
        key TEXT UNIQUE,
        value TEXT,
        created_at DATETIME DEFAULT CURRENT_TIMESTAMP
    )""")
    conn.commit()
    conn.close()


def resolve_workspace_path(path: str) -> str:
    """
    Convert any workspace path (including kernel/physical paths) to its user-visible logical form.
    Handles Windows folder redirection (e.g., Downloads moved from C: to D:).
    """
    if not path: return path
    
    # 1. Basic normalization
    norm = os.path.normpath(path).replace('/', '\\')
    
    # 2. Known folder mapping
    # (USERPROFILE mapping removed as static s->Downloads is preferred for this environment)
    
    # 3. Handle common kernel/physical artifacts (like \s\ or \??\ or device paths)
    # Specifically for the D:\s case mentioned by the user
    if '\\s\\' in norm.lower() or norm.lower().startswith('d:\\s'):
        # Map physical 's' back to 'Downloads'
        logical = norm.replace('\\s\\', '\\Downloads\\').replace('d:\\s', 'd:\\Downloads')
        if os.path.exists(logical):
            return logical

    # 4. Handle redirected 'Downloads' specifically
    # The static mapping above covers the most common cases. 
    # Segment-based logical mapping is handled by dynamic LOGICAL_ROOT mapping in sanitize_path.
    return norm


def sanitize_path(path: str) -> str:
    """Dynamically maps physical paths back to logical paths based on current workspace."""
    if not path or not isinstance(path, str):
        return path
    
    # 1. Use the new resolve_workspace_path for canonicalization
    norm = resolve_workspace_path(path)
    
    # 2. Apply Dynamic Workspace Mappings (LOGICAL_ROOT/PHYSICAL_ROOT)
    if LOGICAL_ROOT and PHYSICAL_ROOT:
        # Normalize for comparison
        n_norm = norm.lower()
        p_norm = PHYSICAL_ROOT.lower()
        if n_norm.startswith(p_norm):
            # Replace physical prefix with logical prefix
            norm = LOGICAL_ROOT + norm[len(PHYSICAL_ROOT):]
    
    # Cleanup double slashes
    while '\\\\' in norm:
        norm = norm.replace('\\\\', '\\')
    
    return norm


def _get_operation_path(path: str) -> str:
    """
    Multi-stage path resolver that prioritizes logical paths.
    1. Try logical path relative to LOGICAL_ROOT.
    2. Try logical path as provided.
    3. Try physical mapping as a fallback.
    4. Try realpath only as a last resort.
    """
    if not path or not isinstance(path, str): return path
    
    # STEP 1: If path is within current workspace, use LOGICAL_ROOT as anchor
    if LOGICAL_ROOT and not os.path.isabs(path):
        candidate = os.path.normpath(os.path.join(LOGICAL_ROOT, path))
        if os.path.exists(candidate):
            return candidate
            
    # Absolute path normalization
    if os.path.isabs(path):
        base = os.path.normpath(path)
    else:
        base = os.path.normpath(os.path.join(CURRENT_WORKSPACE, path))

    # STEP 2: Check existence with logical path directly
    if os.path.exists(base):
        return base

    # STEP 3: Try physical mapping fallback (Downloads -> s)
    candidates = []
    if re.search(r'(?i)[/\\\\]downloads([/\\\\]|$)', base):
        candidates.append(os.path.normpath(re.sub(r'(?i)[/\\\\]downloads([/\\\\]|$)', r'\\s\\', base)))
    if re.search(r'(?i)[/\\\\]s([/\\\\]|$)', base):
        candidates.append(os.path.normpath(re.sub(r'(?i)[/\\\\]s([/\\\\]|$)', r'\\Downloads\\', base)))
    
    # STEP 4: Try realpath (OS-level resolution)
    # CRITICAL: Skip realpath for redirected folders (Downloads) to prevent D:\s leakage
    is_redirected = any(p in base.lower() for p in ['downloads', '\\s\\'])
    if not is_redirected:
        try:
            candidates.append(os.path.normpath(os.path.realpath(base)))
        except:
            pass

    # Test candidates in order
    def is_usable(p):
        try:
            if not os.path.exists(p): return False
            if os.path.isdir(p):
                os.listdir(p) # Verify listable
            else:
                with open(p, 'rb') as f: pass
            return True
        except:
            return False

    for cand in candidates:
        if is_usable(cand):
            return cand
            
    return base


def diagnose_path(path: str) -> Dict[str, Any]:
    """Provides deep diagnostic info about path visibility for debugging."""
    if not path: return {"error": "No path provided"}
    full = os.path.normpath(os.path.join(CURRENT_WORKSPACE, path)) if not os.path.isabs(path) else os.path.normpath(path)
    
    results = {
        "input_path": path,
        "absolute_normalized": full,
        "exists": os.path.exists(full),
        "is_dir": os.path.isdir(full) if os.path.exists(full) else None,
        "realpath": os.path.realpath(full),
        "working_path": _get_operation_path(path),
        "current_workspace": CURRENT_WORKSPACE
    }
    
    # Junction/Symlink detection (Windows specific)
    try:
        import ctypes
        FILE_ATTRIBUTE_REPARSE_POINT = 0x400
        attrs = ctypes.windll.kernel32.GetFileAttributesW(full)
        results["is_junction"] = bool(attrs & FILE_ATTRIBUTE_REPARSE_POINT) if attrs != 0xFFFFFFFF else False
        if results["is_junction"]:
            import subprocess
            res = subprocess.run(['fsutil', 'reparsepoint', 'query', full], capture_output=True, text=True, shell=True)
            results["junction_info"] = res.stdout[:500] if res.returncode == 0 else "fsutil failed"
    except:
        results["is_junction"] = False
    
    # Try shell visibility
    try:
        import subprocess
        res = subprocess.run(f'dir "{full}"', shell=True, capture_output=True, text=True)
        results["shell_visibility"] = "Visible" if res.returncode == 0 else f"Hidden (Code {res.returncode})"
        if res.stdout:
            results["shell_output_snippet"] = res.stdout[:500]
    except Exception as e:
        results["shell_error"] = str(e)
        
    return results


def set_workspace(path: str) -> Dict[str, Any]:
    global CURRENT_WORKSPACE, LOGICAL_ROOT, PHYSICAL_ROOT
    
    # 1. Canonicalize the workspace path to its logical form
    logical_workspace = resolve_workspace_path(path)
    
    # 2. Check if logical exists, otherwise fallback to physical
    if os.path.exists(logical_workspace):
        CURRENT_WORKSPACE = logical_workspace
    elif os.path.exists(path):
        CURRENT_WORKSPACE = os.path.normpath(path)
    else:
        return {"error": f"Path not found: {path}"}
        
    # 3. Establish Junction Mappings for this session
    # We use normpath and realpath to detect the physical underlying folder
    LOGICAL_ROOT = CURRENT_WORKSPACE
    try:
        # Detect if this logical path points to a different physical path
        physical = os.path.normpath(os.path.realpath(CURRENT_WORKSPACE))
        if physical.lower() != CURRENT_WORKSPACE.lower():
            PHYSICAL_ROOT = physical
        else:
            # Check for the common 's' junction on D:
            if 'downloads' in CURRENT_WORKSPACE.lower():
                PHYSICAL_ROOT = CURRENT_WORKSPACE.lower().replace('\\downloads', '\\s')
            else:
                PHYSICAL_ROOT = CURRENT_WORKSPACE
    except:
        PHYSICAL_ROOT = CURRENT_WORKSPACE
        
    _init_db()
    return {
        "status": "success", 
        "workspace": CURRENT_WORKSPACE,
        "mapping": f"{PHYSICAL_ROOT} -> {LOGICAL_ROOT}" if LOGICAL_ROOT != PHYSICAL_ROOT else "Direct"
    }
    
    # Fallback: maybe the user provided a physical path? Try to map it back
    if "\\s\\" in normalized_path.lower():
        import re
        # Convert physical path back to user-friendly path
        user_path = re.sub(r'\\s\\', r'\\Downloads\\', normalized_path, flags=re.IGNORECASE)
        if os.path.exists(user_path):
            CURRENT_WORKSPACE = user_path
            _init_db()
            return {"status": "success", "workspace": CURRENT_WORKSPACE, "note": "Mapped from physical to user path"}
    
    return {"error": f"'{path}' is not a directory"}


def requires_approval(tool_name: str, args: Dict[str, Any]) -> bool:
    if tool_name in {"run_command", "run_shell_command", "start_background_command"}:
        cmd = args.get("command", "").lower()
        return not any(cmd.startswith(safe) for safe in SAFE_COMMANDS)
        
    if tool_name in {"write_file", "append_file"}:
        path = args.get("path")
        if path:
            full = _get_operation_path(path)
            # Require approval ONLY if the file already exists
            if os.path.exists(full):
                return True
        return False
        
    return tool_name in {
        "edit_file", "move_file", "copy_file", "delete_file"
    }


def get_diff(path: str, new_content: str) -> str:
    if not isinstance(path, str):
        return ""
    full_path = os.path.isabs(path) and path or os.path.join(CURRENT_WORKSPACE, path)
    old = ""
    if os.path.exists(full_path):
        try:
            with open(full_path, "r", encoding="utf-8") as f:
                old = f.read()
        except Exception:
            pass
    diff = difflib.unified_diff(
        old.splitlines(keepends=True),
        new_content.splitlines(keepends=True),
        fromfile="original",
        tofile="proposed"
    )
    return "".join(diff)


# ---------------------------------------------------------------------------
# File system tools
# ---------------------------------------------------------------------------
def _read_file_content(full: str, ext: str) -> Dict[str, Any]:
    """Unified helper to extract content from various file types."""
    # Image handling (Vision)
    IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".ico", ".svg", ".tiff"}
    if ext in IMAGE_EXTS:
        with open(full, "rb") as f:
            b64 = base64.b64encode(f.read()).decode()
        return {"content": b64, "is_image": True, "ext": ext.lstrip(".")}

    # PDF handling
    if ext == ".pdf":
        try:
            import pypdf
            reader = pypdf.PdfReader(full)
            text = "\n".join([page.extract_text() or "" for page in reader.pages])
            return {"content": text, "is_pdf": True, "page_count": len(reader.pages)}
        except Exception as e:
            return {"error": f"PDF Error: {str(e)}"}

    # Word Documents
    if ext in {".docx", ".doc"}:
        try:
            import docx
            doc = docx.Document(full)
            text = "\n".join([p.text for p in doc.paragraphs])
            return {"content": text, "is_doc": True}
        except Exception as e:
            return {"error": f"Word Error: {str(e)}"}

    # PowerPoint
    if ext in {".pptx", ".ppt"}:
        try:
            from pptx import Presentation
            prs = Presentation(full)
            text = ""
            for i, slide in enumerate(prs.slides):
                text += f"--- Slide {i+1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return {"content": text, "is_pptx": True}
        except Exception as e:
            return {"error": f"PowerPoint Error: {str(e)}"}

    # EPUB
    if ext == ".epub":
        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup
            book = epub.read_epub(full)
            text = ""
            for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
                text += BeautifulSoup(item.get_content(), 'html.parser').get_text() + "\n"
            return {"content": text, "is_epub": True}
        except Exception as e:
            return {"error": f"EPUB Error: {str(e)}"}

    # Data Files (Pandas)
    if ext in {".xlsx", ".xls", ".csv", ".parquet"}:
        try:
            import pandas as pd
            if ext == ".parquet":
                df = pd.read_parquet(full)
            elif ext in {".xlsx", ".xls"}:
                df = pd.read_excel(full)
            else:
                df = pd.read_csv(full)
            return {"content": df.to_csv(index=False), "is_data": True, "shape": df.shape}
        except Exception as e:
            if ext != ".csv": return {"error": f"Data Error: {str(e)}"}
            
    return {"is_text": True}


def read_file(path: str, start_line: Optional[int] = None, end_line: Optional[int] = None) -> Dict[str, Any]:
    full = _get_operation_path(path)
    if not os.path.exists(full):
        return {"error": f"File not found: {path}"}
    ext = os.path.splitext(full)[1].lower()
    
    # Try specialized extraction
    special = _read_file_content(full, ext)
    if "error" in special: return special
    if not special.get("is_text"):
        # For images, we need to add the path back for the agent_brain to use
        if special.get("is_image"):
            special["path"] = path
        return special

    # Plain Text Logic

    mtime = os.path.getmtime(full)
    if full in _file_cache and _file_cache[full][0] == mtime:
        content = _file_cache[full][1]
    else:
        for enc in ("utf-8-sig", "utf-8", "utf-16", "latin-1"):
            try:
                with open(full, "r", encoding=enc) as f:
                    content = f.read()
                break
            except UnicodeError:
                continue
        else:
            with open(full, "rb") as f:
                content = f.read().decode("utf-8", errors="replace")
        _file_cache[full] = (mtime, content)

    lines = content.splitlines()
    total = len(lines)
    s = max(0, (start_line - 1) if start_line else 0)
    e = min(total, end_line if end_line else total)
    MAX_LINES = 1000
    note = ""
    if e - s > MAX_LINES:
        e = s + MAX_LINES
        note = f"⚠️ Truncated to {e - s} lines (lines {s+1}–{e} of {total}). Use start_line/end_line to paginate."
    elif start_line or end_line:
        note = f"Showing lines {s+1}–{e} of {total}."
    chunk = "\n".join(lines[s:e])
    return {"content": chunk, "range": [s+1, e], "total_lines": total, "note": note}


def read_files(paths: List[str]) -> Dict[str, Any]:
    """
    Read multiple small/medium files at once. 
    Agent Brain dynamically decides which files to batch based on size.
    """
    if not isinstance(paths, list):
        return {"error": "paths must be a list of strings"}
    
    results = {}
    errors = {}
    total_chars = 0
    MAX_BATCH_CHARS = 100000 # ~100KB limit for the entire batch
    MAX_SINGLE_FILE_CHARS = 30000 # Skip individual files larger than 30KB in batch mode
    
    for path in paths:
        try:
            full = _get_operation_path(path)
            if not os.path.exists(full):
                errors[path] = "Not found"
                continue
                
            ext = os.path.splitext(full)[1].lower()
            res = _read_file_content(full, ext)
            
            if "error" in res:
                errors[path] = res["error"]
                continue
                
            content = res.get("content", "")
            
            # Special handling for images in batch
            if res.get("is_image"):
                content = f"IMAGE:{content}"
            elif res.get("is_text"):
                # Load text content with encoding logic
                # (We still need the encoding loop for unknown text files)
                mtime = os.path.getmtime(full)
                if full in _file_cache and _file_cache[full][0] == mtime:
                    content = _file_cache[full][1]
                else:
                    for enc in ("utf-8-sig", "utf-8", "utf-16", "latin-1"):
                        try:
                            with open(full, "r", encoding=enc) as f:
                                content = f.read()
                            break
                        except UnicodeError:
                            continue
                    else:
                        with open(full, "rb") as f:
                            content = f.read().decode("utf-8", errors="replace")
                    _file_cache[full] = (mtime, content)
            
            if len(content) > MAX_SINGLE_FILE_CHARS:
                errors[path] = f"File too large for batch ({len(content)} chars). Use read_file individually."
                continue

            if total_chars + len(content) > MAX_BATCH_CHARS:
                errors[path] = "Batch capacity reached"
                break
                
            results[path] = content
            total_chars += len(content)
            
        except Exception as e:
            errors[path] = str(e)
            
    return {
        "files": results, 
        "errors": errors if errors else None,
        "note": f"Read {len(results)} files. Total size: {total_chars} chars."
    }


def write_file(path: str, content: str = "", lines: Optional[List[str]] = None,
               content_base64: str = "", lines_base64: Optional[List[str]] = None) -> Dict[str, Any]:
    # ENFORCEMENT: For code files, require lines, base64, or non-empty content (raw block)
    is_code_file = path.lower().endswith(('.py', '.js', '.ts', '.html', '.css', '.json', '.yaml', '.yml', '.md', '.sh', '.bat'))
    if is_code_file and not (content_base64 or lines_base64 or lines is not None or content):
        return {
            "error": f"Code file '{path}' requires 'lines' (array) or base64 encoding to preserve indentation. Please use 'lines' (recommended for Unicode) or 'lines_base64'.",
            "hint": "Raw multiline 'content' string strips leading spaces in JSON. Use 'lines'."
        }
    if content_base64:
        content = _decode_base64(content_base64)
    elif lines_base64 is not None:
        lines = [_decode_base64(l) for l in lines_base64]
        content = "\n".join(lines)
    elif lines is not None:
        content = "\n".join(lines)

    # Built-in syntax check validation guards before writing
    if path.lower().endswith(".py"):
        syntax_err = check_python_syntax(content)
        if syntax_err:
            return {
                "error": "Python syntax validation failed before writing.",
                "details": syntax_err,
                "hint": "Please review your code structure, indentation, or statements, and rewrite a valid script."
            }
    elif path.lower().endswith(".json"):
        syntax_err = check_json_syntax(content)
        if syntax_err:
            return {
                "error": "JSON validation failed before writing.",
                "details": syntax_err,
                "hint": "Please verify your braces, commas, and double-quotes to ensure standard JSON compliance."
            }

    full = _get_operation_path(path)
    try:
        os.makedirs(os.path.dirname(full), exist_ok=True)
        with open(full, "w", encoding="utf-8") as f:
            f.write(content)
        _file_cache.pop(full, None)
        return {"status": "success", "path": sanitize_path(full)}
    except Exception as e:
        return {"error": str(e)}


def edit_file(path: str, target: str = "", replacement: str = "", 
              target_lines: Optional[List[str]] = None, replacement_lines: Optional[List[str]] = None,
              target_base64: str = "", replacement_base64: str = "",
              target_lines_base64: Optional[List[str]] = None, replacement_lines_base64: Optional[List[str]] = None) -> Dict[str, Any]:
    # ENFORCEMENT: For code files, require lines or base64 to preserve indentation
    is_code_file = path.lower().endswith(('.py', '.js', '.ts', '.html', '.css', '.json', '.yaml', '.yml', '.md', '.sh', '.bat'))
    if is_code_file and not (target_base64 or replacement_base64 or target_lines_base64 or replacement_lines_base64 or target_lines is not None or replacement_lines is not None):
         return {
            "error": f"Editing code file '{path}' requires 'lines' (array) or base64 encoding to preserve indentation. Please use 'target_lines' and 'replacement_lines'.",
            "hint": "Raw multiline 'content' string strips leading spaces in JSON. Use 'lines'."
        }
    if target_base64:
        target = _decode_base64(target_base64)
    if replacement_base64:
        replacement = _decode_base64(replacement_base64)
    if target_lines_base64:
        target_lines = [_decode_base64(l) for l in target_lines_base64]
        target = "\n".join(target_lines)
    if replacement_lines_base64:
        replacement_lines = [_decode_base64(l) for l in replacement_lines_base64]
        replacement = "\n".join(replacement_lines)
        
    if target_lines is not None:
        target = "\n".join(target_lines)
    if replacement_lines is not None:
        replacement = "\n".join(replacement_lines)
    full = _get_operation_path(path)
    if not os.path.exists(full):
        return {"error": f"File '{path}' not found"}
    with open(full, "r", encoding="utf-8") as f:
        original = f.read()

    new_content = None
    match_type = None

    # Stage 1: exact match
    if original.count(target) == 1:
        new_content = original.replace(target, replacement, 1)
        match_type = "exact"

    # Stage 2: whitespace-agnostic regex
    if new_content is None:
        target_lines = [l.strip() for l in target.splitlines() if l.strip()]
        if not target_lines:
            return {"error": "Target is empty or whitespace only."}
        escaped = [re.escape(l) for l in target_lines]
        pattern = r'[ \t]*' + r'\s+'.join(escaped) + r'[ \t]*'
        matches = list(re.finditer(pattern, original, re.MULTILINE))
        if len(matches) == 1:
            new_content = original.replace(matches[0].group(0), replacement, 1)
            match_type = "fuzzy_regex"
        elif len(matches) > 1:
            return {"error": f"Found {len(matches)} fuzzy matches. Provide more context."}

    # Stage 3: first & last line anchor
    if new_content is None:
        first_line, last_line = target_lines[0], target_lines[-1]
        candidate = None
        for i, line in enumerate(original.splitlines()):
            if line.strip() == first_line:
                for j in range(i, min(i + len(target_lines) + 10, len(original.splitlines()))):
                    if original.splitlines()[j].strip() == last_line:
                        if candidate is None:
                            candidate = (i, j)
                        else:
                            return {"error": "Multiple anchor matches found."}
        if candidate:
            start_idx, end_idx = candidate
            lines = original.splitlines()
            new_lines = lines[:start_idx] + [replacement] + lines[end_idx+1:]
            new_content = "\n".join(new_lines)
            if original.endswith("\n") and not new_content.endswith("\n"):
                new_content += "\n"
            match_type = "anchor"

    if new_content is None:
        return {"error": "Target not found."}

    # Built-in syntax check guards for the edited file
    if path.lower().endswith(".py"):
        syntax_err = check_python_syntax(new_content)
        if syntax_err:
            return {
                "error": "Python syntax validation failed after applying this edit. Edit aborted.",
                "details": syntax_err,
                "hint": "Please review the replacement content, indentation, and structure to ensure it remains a valid script."
            }
    elif path.lower().endswith(".json"):
        syntax_err = check_json_syntax(new_content)
        if syntax_err:
            return {
                "error": "JSON validation failed after applying this edit. Edit aborted.",
                "details": syntax_err,
                "hint": "Please check your braces, commas, or quotes in the replacement content to keep standard JSON compatibility."
            }

    # Write the validated new content
    try:
        with open(full, "w", encoding="utf-8") as f:
            f.write(new_content)
        _file_cache.pop(full, None)
        return {"status": "success", "match_type": match_type}
    except Exception as e:
        return {"error": str(e)}


def append_file(path: str, content: str = "", lines: Optional[List[str]] = None,
                content_base64: str = "", lines_base64: Optional[List[str]] = None) -> Dict[str, Any]:
    if content_base64:
        content = _decode_base64(content_base64)
    elif lines_base64 is not None:
        lines = [_decode_base64(l) for l in lines_base64]
        content = "\n".join(lines)
    elif lines is not None:
        content = "\n".join(lines)
    full = _get_operation_path(path)
    try:
        os.makedirs(os.path.dirname(full), exist_ok=True)
        with open(full, "a", encoding="utf-8") as f:
            f.write(content)
        return {"status": "success", "path": sanitize_path(full)}
    except Exception as e:
        return {"error": str(e)}


def create_folder(path: str) -> Dict[str, Any]:
    full = _get_operation_path(path)
    os.makedirs(full, exist_ok=True)
    return {"status": "success", "path": sanitize_path(full)}


def delete_file(path: str) -> Dict[str, Any]:
    full = _get_operation_path(path)
    if not os.path.exists(full):
        return {"error": f"'{path}' not found"}
    if os.path.isdir(full):
        shutil.rmtree(full)
    else:
        os.remove(full)
    return {"status": "success", "path": sanitize_path(full)}


def move_file(source: str, destination: str) -> Dict[str, Any]:
    src = _get_operation_path(source)
    dst = _get_operation_path(destination)
    os.makedirs(os.path.dirname(dst), exist_ok=True)
    shutil.move(src, dst)
    return {"status": "success", "from": sanitize_path(src), "to": sanitize_path(dst)}


def copy_file(source: str, destination: str) -> Dict[str, Any]:
    src = _get_operation_path(source)
    dst = _get_operation_path(destination)
    os.makedirs(os.path.dirname(dst), exist_ok=True)
    if os.path.isdir(src):
        shutil.copytree(src, dst, dirs_exist_ok=True)
    else:
        shutil.copy2(src, dst)
    return {"status": "success", "from": sanitize_path(src), "to": sanitize_path(dst)}


def get_file_info(path: str) -> Dict[str, Any]:
    full = _get_operation_path(path)
    if not os.path.exists(full):
        return {"error": f"'{path}' not found"}
    stat = os.stat(full)
    logical_path = os.path.normpath(full)
    physical_path = os.path.realpath(full)
    return {
        "path": path,
        "type": "folder" if os.path.isdir(full) else "file",
        "size_bytes": stat.st_size,
        "modified_at": datetime.datetime.fromtimestamp(stat.st_mtime).isoformat(),
        "resolved_physical_path": physical_path,
        "is_logical_path": logical_path.lower() != physical_path.lower(),
        "is_mapped": logical_path.lower() != physical_path.lower()
    }


def resolve_path(path: str) -> Dict[str, Any]:
    full = _get_operation_path(path)
    logical_path = sanitize_path(full)
    physical_path = os.path.realpath(full)
    return {
        "logical_path": logical_path,
        "physical_path": physical_path,
        "is_mapped": logical_path.lower() != physical_path.lower()
    }


def list_files(path: str = ".") -> Dict[str, Any]:
    full = _get_operation_path(path)
    try:
        entries = os.listdir(full)
        valid_entries = []
        dead_junctions = []
        for entry in entries:
            entry_path = os.path.join(full, entry)
            # Skip dead junctions (dir entry exists but target is missing)
            if os.path.isdir(entry_path) and not os.path.exists(entry_path):
                dead_junctions.append(entry)
                continue
            valid_entries.append(entry)
            
        res = {"files": valid_entries, "path": sanitize_path(full)}
        if dead_junctions:
            res["note"] = f"Filtered out {len(dead_junctions)} dead junctions: {dead_junctions}"
        return res
    except Exception as e:
        return {"error": str(e)}


def verify_junction(path: str) -> Dict[str, Any]:
    """Check if path is a junction point and verify its target exists."""
    full = _get_operation_path(path)
    results = {"path": sanitize_path(full), "exists": os.path.exists(full)}
    try:
        import ctypes
        attrs = ctypes.windll.kernel32.GetFileAttributesW(full)
        results["is_junction"] = bool(attrs & 0x400) if attrs != 0xFFFFFFFF else False
        if results["is_junction"]:
            import subprocess
            res = subprocess.run(['fsutil', 'reparsepoint', 'query', full], capture_output=True, text=True, shell=True)
            if res.returncode == 0:
                match = re.search(r'Substitute Name:.*?([A-Za-z]:\\[^\n\r]+)', res.stdout)
                target = match.group(1).replace('\\??\\', '') if match else "Unknown"
                results["target"] = target
                results["target_exists"] = os.path.exists(target)
    except Exception as e:
        results["error"] = str(e)
    return results


def resolve_junction(path: str) -> Dict[str, Any]:
    """Follow a junction point to its physical target using Windows API."""
    full = _get_operation_path(path)
    try:
        import ctypes
        kernel32 = ctypes.windll.kernel32
        # FILE_FLAG_BACKUP_SEMANTICS = 0x02000000, OPEN_EXISTING = 3, GENERIC_READ = 0x80000000
        handle = kernel32.CreateFileW(full, 0x80000000, 1, None, 3, 0x02000000, None)
        if handle == -1: return {"error": "Failed to open handle to junction"}
        buf = ctypes.create_unicode_buffer(1024)
        kernel32.GetFinalPathNameByHandleW(handle, buf, 1024, 0)
        kernel32.CloseHandle(handle)
        physical = buf.value.replace('\\\\?\\', '')
        return {"original": sanitize_path(full), "physical": physical, "exists": os.path.exists(physical)}
    except Exception as e:
        return {"error": str(e)}


def get_file_tree(startpath: Optional[str] = None) -> Any:
    # Use sanitized workspace for relative calculations
    if startpath is None:
        startpath = CURRENT_WORKSPACE
    root = sanitize_path(os.path.normpath(CURRENT_WORKSPACE))
    budget = {"remaining": MAX_TREE_NODES}

    def walk(p: str) -> Optional[Dict]:
        if budget["remaining"] <= 0:
            return None
        name = os.path.basename(p) or os.path.basename(p.rstrip(os.sep))
        if name in IGNORED_TREE_DIRS:
            return None
        
        # Aggressively sanitize path to logical form before calculating relative path
        p_norm = sanitize_path(p)
        rel = "." if p_norm == root else os.path.relpath(p_norm, root)
        
        node = {"name": name, "path": rel}
        budget["remaining"] -= 1
        if os.path.isdir(p):
            node["type"] = "folder"
            if p_norm != root and name in TRIMMED_TREE_DIRS:
                node["children"] = []
                node["trimmed"] = True
                return node
            children = []
            try:
                for entry in sorted(os.listdir(p)):
                    if budget["remaining"] <= 0:
                        node["trimmed"] = True
                        break
                    child = walk(os.path.join(p, entry))
                    if child:
                        children.append(child)
            except (PermissionError, OSError):
                pass
            node["children"] = children
        else:
            node["type"] = "file"
        return node

    return walk(startpath)


def grep_search(pattern: str, path: str = ".", case_insensitive: bool = True) -> Dict[str, Any]:
    root = _get_operation_path(path)
    flags = re.IGNORECASE if case_insensitive else 0
    try:
        regex = re.compile(pattern, flags)
    except re.error as e:
        return {"error": f"Invalid regex: {e}"}
    results = []
    max_results = 100
    for dirpath, dirnames, filenames in os.walk(root):
        dirnames[:] = [d for d in dirnames if d not in IGNORED_TREE_DIRS and d not in TRIMMED_TREE_DIRS]
        for fname in filenames:
            if os.path.splitext(fname)[1].lower() in {".png", ".jpg", ".exe", ".dll", ".pyc", ".o", ".bin"}:
                continue
            full = os.path.join(dirpath, fname)
            try:
                with open(full, "r", encoding="utf-8", errors="ignore") as f:
                    for i, line in enumerate(f, 1):
                        if regex.search(line):
                            results.append({"file": os.path.relpath(sanitize_path(full), CURRENT_WORKSPACE), "line": i, "content": line.strip()})
                            if len(results) >= max_results:
                                return {"results": results, "note": f"Capped at {max_results}"}
            except Exception:
                continue
    return {"results": results}


def find_files(pattern: str) -> Dict[str, Any]:
    import fnmatch
    results = []
    max_results = 200
    for dirpath, dirnames, filenames in os.walk(_get_operation_path(".")):
        dirnames[:] = [d for d in dirnames if d not in IGNORED_TREE_DIRS and d not in TRIMMED_TREE_DIRS]
        for fname in filenames:
            if fnmatch.fnmatch(fname, pattern) or pattern.lower() in fname.lower():
                results.append(os.path.relpath(sanitize_path(os.path.join(dirpath, fname)), CURRENT_WORKSPACE))
                if len(results) >= max_results:
                    break
        if len(results) >= max_results:
            break
    return {"files": results, "note": f"Found {len(results)} matches"}


def reveal_in_os(path: str) -> Dict[str, Any]:
    full = _get_operation_path(path)
    if not os.path.exists(full):
        return {"error": f"'{path}' not found"}
    if os.name == "nt":
        subprocess.Popen(f'explorer /select,"{full}"')
    else:
        subprocess.Popen(["open", "-R", full])
    return {"status": "success"}


# ---------------------------------------------------------------------------
# Shell command execution (with live streaming)
# ---------------------------------------------------------------------------
async def run_command_async(command: str, callback: Optional[Callable] = None) -> Dict[str, Any]:
    if callback:
        await callback({"type": "direct_terminal_result",
                         "html": f"<div class='terminal-cmd'>▶ {command}</div>",
                         "agent_controlled": True})
    loop = asyncio.get_running_loop()

    def _run():
        try:
            proc = subprocess.Popen(command, shell=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE,
                                   cwd=CURRENT_WORKSPACE, text=False)
            temp_id = str(uuid.uuid4())[:8]
            active_processes[temp_id] = {"process": proc, "command": command, "buffer": []}

            def stream(pipe, is_stderr: bool):
                while True:
                    chunk = pipe.read(1024)
                    if not chunk:
                        break
                    decoded = chunk.decode("utf-8", errors="replace")
                    active_processes[temp_id]["buffer"].append(decoded)
                    if is_stderr:
                        if callback:
                            asyncio.run_coroutine_threadsafe(
                                callback({"type": "direct_terminal_result", "stderr": decoded, "agent_controlled": True}),
                                loop
                            )
                    else:
                        if callback:
                            asyncio.run_coroutine_threadsafe(
                                callback({"type": "direct_terminal_result", "stdout": decoded, "agent_controlled": True}),
                                loop
                            )

            t1 = threading.Thread(target=stream, args=(proc.stdout, False), daemon=True)
            t2 = threading.Thread(target=stream, args=(proc.stderr, True), daemon=True)
            t1.start(); t2.start()

            try:
                proc.wait(timeout=15)
            except subprocess.TimeoutExpired:
                proc.kill()
                t1.join(); t2.join()
                active_processes.pop(temp_id, None)
                return {"error": "Command timed out after 15 seconds. Use start_background_command for long‑running tasks."}

            t1.join(); t2.join()
            raw_stdout = "".join(active_processes[temp_id]["buffer"])
            active_processes.pop(temp_id, None)
            truncated_stdout = truncate_stdout(raw_stdout, max_lines=300)
            return {"stdout": truncated_stdout, "stderr": "", "code": proc.returncode}
        except Exception as e:
            return {"error": str(e)}

    result = await loop.run_in_executor(None, _run)
    if callback:
        await callback({"type": "direct_terminal_result",
                         "html": f"<div class='terminal-exit'>Exit code: {result.get('code', 'N/A')}</div>",
                         "agent_controlled": True})
    return result


async def start_background_command(command: str, callback: Optional[Callable] = None) -> Dict[str, Any]:
    pid = str(uuid.uuid4())[:8]
    if callback:
        await callback({"type": "direct_terminal_result",
                         "html": f"<div class='terminal-cmd'>▶ Background [{pid}]: {command}</div>",
                         "agent_controlled": True})
    loop = asyncio.get_running_loop()

    def _run():
        proc = subprocess.Popen(command, shell=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE,
                               cwd=CURRENT_WORKSPACE, bufsize=1, text=False)
        buf = []
        active_processes[pid] = {"process": proc, "command": command, "buffer": buf}

        def reader(pipe, is_stderr):
            for line in pipe:
                decoded = line.decode("utf-8", errors="replace")
                buf.append(decoded)
                if callback:
                    asyncio.run_coroutine_threadsafe(
                        callback({"type": "direct_terminal_result", "stderr" if is_stderr else "stdout": decoded, "agent_controlled": True}),
                        loop
                    )

        threading.Thread(target=reader, args=(proc.stdout, False), daemon=True).start()
        threading.Thread(target=reader, args=(proc.stderr, True), daemon=True).start()
        return {"status": "started", "process_id": pid, "command": command}

    return await loop.run_in_executor(None, _run)


def read_process_output(process_id: str) -> Dict[str, Any]:
    if process_id not in active_processes:
        return {"error": "Unknown process ID"}
    info = active_processes[process_id]
    proc = info["process"]
    out = "".join(info["buffer"])
    info["buffer"].clear()
    truncated_out = truncate_stdout(out, max_lines=300)
    status = "running" if proc.poll() is None else f"exited with code {proc.poll()}"
    if status != "running":
        del active_processes[process_id]
    return {"process_id": process_id, "status": status, "output": truncated_out or "(no new output)"}


def kill_process(process_id: str) -> Dict[str, Any]:
    if process_id not in active_processes:
        return {"error": "Unknown process ID"}
    info = active_processes.pop(process_id)
    proc = info["process"]
    if proc.poll() is None:
        proc.kill()
        return {"status": "killed", "process_id": process_id}
    return {"status": "already exited", "process_id": process_id}


def list_running_processes() -> Dict[str, Any]:
    running = []
    for pid, info in list(active_processes.items()):
        if info["process"].poll() is None:
            running.append({"process_id": pid, "command": info["command"], "status": "running"})
        else:
            del active_processes[pid]
    return {"running_processes": running}


# ---------------------------------------------------------------------------
# Web tools
# ---------------------------------------------------------------------------
async def search_web(query: str) -> Dict[str, Any]:
    try:
        async with httpx.AsyncClient(timeout=10) as client:
            resp = await client.get(
                f"https://lite.duckduckgo.com/lite/?q={url_quote(query)}",
                headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"}
            )
            html = resp.text
        results = []
        for m in re.finditer(r'<a[^>]*class="result-link"[^>]*href="([^"]+)"[^>]*>([^<]+)</a>.*?<td class="result-snippet">([^<]+)</td>', html, re.DOTALL):
            link, title, snippet = m.groups()
            results.append({"title": title.strip(), "link": link if link.startswith("http") else f"https:{link}", "snippet": snippet.strip()})
        return {"results": results[:8]}
    except Exception as e:
        return {"error": str(e)}


async def read_url(url: str) -> Dict[str, Any]:
    try:
        async with httpx.AsyncClient(timeout=15, follow_redirects=True) as client:
            resp = await client.get(url, headers={"User-Agent": "Mozilla/5.0 ..."})
            html = resp.text
        # Simple HTML → text
        html = re.sub(r'<(script|style)[^>]*>.*?</\1>', '', html, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r'<(p|br|div|li|h\d)[^>]*>', '\n', html, flags=re.IGNORECASE)
        text = re.sub(r'<[^>]+>', '', text)
        import html as html_mod
        text = html_mod.unescape(text)
        text = re.sub(r'\n\s*\n', '\n\n', text).strip()
        return {"content": text[:4000] + ("..." if len(text) > 4000 else ""), "url": url}
    except Exception as e:
        return {"error": str(e)}


# ---------------------------------------------------------------------------
# Memory & knowledge base
# ---------------------------------------------------------------------------
def store_memory(key: str, value: str) -> Dict[str, Any]:
    try:
        conn = sqlite3.connect(os.path.join(CURRENT_WORKSPACE, "memories.db"))
        conn.execute("INSERT OR REPLACE INTO memories (key, value) VALUES (?, ?)", (key, value))
        conn.commit()
        conn.close()
        return {"status": "success"}
    except Exception as e:
        return {"error": str(e)}


def recall_memory(query: str) -> Dict[str, Any]:
    try:
        conn = sqlite3.connect(os.path.join(CURRENT_WORKSPACE, "memories.db"))
        cur = conn.execute("SELECT key, value FROM memories WHERE key LIKE ? OR value LIKE ?", (f"%{query}%", f"%{query}%"))
        results = [{"key": r[0], "value": r[1]} for r in cur.fetchall()]
        conn.close()
        return {"results": results}
    except Exception as e:
        return {"error": str(e)}


# Knowledge base functions are imported from knowledge_base.py; they are available as
# `learn_pattern`, `recall_pattern`, `update_pattern_success`, `get_knowledge_stats`.


# ---------------------------------------------------------------------------
# Git helpers
# ---------------------------------------------------------------------------
def git_commit(message: str, files: Optional[List[str]] = None) -> Dict[str, Any]:
    try:
        import git
        repo = git.Repo(CURRENT_WORKSPACE)
        if files:
            repo.index.add(files)
        else:
            repo.index.add("*")
        repo.index.commit(message)
        return {"commit": str(repo.head.commit.hexsha)}
    except Exception as e:
        return {"error": str(e)}


def git_diff(path: Optional[str] = None) -> Dict[str, Any]:
    try:
        cmd = ["git", "diff"]
        if path:
            full = _get_operation_path(path)
            cmd += ["--", full]
        p = subprocess.run(cmd, capture_output=True, text=True, cwd=CURRENT_WORKSPACE)
        return {"diff": p.stdout or "(No changes)", "stderr": p.stderr, "code": p.returncode}
    except Exception as e:
        return {"error": str(e)}


def run_tests(path: str = ".", pattern: Optional[str] = None) -> Dict[str, Any]:
    cmd = ["pytest", path, "-v", "--tb=short"]
    if pattern:
        cmd.extend(["-k", pattern])
    result = subprocess.run(cmd, capture_output=True, text=True, cwd=CURRENT_WORKSPACE)
    return {"passed": result.returncode == 0, "stdout": result.stdout, "stderr": result.stderr}


def get_definition(file: str, line: int, column: int) -> Dict[str, Any]:
    try:
        import jedi
        full = os.path.isabs(file) and file or os.path.join(CURRENT_WORKSPACE, file)
        script = jedi.Script(path=full)
        defs = script.goto(line, column) if hasattr(script, "goto") else script.goto_definitions(line, column)
        return {"definitions": [{
            "name": d.name,
            "module_path": str(d.module_path),
            "line": d.line,
            "column": d.column,
            "description": d.description
        } for d in defs]}
    except Exception as e:
        return {"error": str(e)}


# =============================================================================
# EXTENDED TOOLS (PRIORITIES 1-3)
# =============================================================================

def get_env_var(name: str) -> Dict[str, Any]:
    """Get environment variable value."""
    value = os.environ.get(name)
    return {"name": name, "value": value, "exists": value is not None}


def set_env_var(name: str, value: str, permanent: bool = False) -> Dict[str, Any]:
    """Set environment variable."""
    os.environ[name] = value
    
    if permanent and sys.platform == "win32":
        try:
            import ctypes
            ctypes.windll.user32.SendMessageTimeoutW(0xFFFF, 0x001A, 0, "Environment", 0x0002, 5000, None)
        except:
            pass
    
    return {"status": "success", "name": name, "value": value}


def get_system_info() -> Dict[str, Any]:
    """Get comprehensive system information."""
    info = {
        "os": platform.system(),
        "os_version": platform.version(),
        "os_release": platform.release(),
        "machine": platform.machine(),
        "processor": platform.processor(),
        "python_version": sys.version,
        "cpu_count": os.cpu_count(),
        "hostname": platform.node(),
    }
    
    if HAVE_PSUTIL:
        import psutil
        mem = psutil.virtual_memory()
        info["ram_total_gb"] = round(mem.total / (1024**3), 2)
        info["ram_available_gb"] = round(mem.available / (1024**3), 2)
        info["ram_percent_used"] = mem.percent
        
        disk = psutil.disk_usage('/')
        info["disk_total_gb"] = round(disk.total / (1024**3), 2)
        info["disk_free_gb"] = round(disk.free / (1024**3), 2)
        info["disk_percent_used"] = disk.percent
    
    return info


def is_admin() -> Dict[str, Any]:
    """Check for admin/root privileges."""
    if sys.platform == "win32":
        try:
            import ctypes
            is_admin = ctypes.windll.shell32.IsUserAnAdmin() != 0
        except:
            is_admin = False
    else:
        is_admin = os.getuid() == 0 if hasattr(os, 'getuid') else False
    
    return {"is_admin": is_admin}


def download_file(url: str, dest: str, resume: bool = False, timeout: int = 300) -> Dict[str, Any]:
    """Download a file with resume support."""
    dest_path = Path(dest)
    dest_path.parent.mkdir(parents=True, exist_ok=True)
    
    existing_size = dest_path.stat().st_size if resume and dest_path.exists() else 0
    headers = {'Range': f'bytes={existing_size}-'} if existing_size > 0 else {}
    
    try:
        import urllib.request
        req = urllib.request.Request(url, headers=headers)
        with urllib.request.urlopen(req, timeout=timeout) as response:
            mode = 'ab' if existing_size > 0 else 'wb'
            with open(dest_path, mode) as f:
                while True:
                    chunk = response.read(8192)
                    if not chunk:
                        break
                    f.write(chunk)
        
        return {
            "status": "success",
            "path": str(dest_path),
            "size_bytes": dest_path.stat().st_size,
            "resumed": existing_size > 0
        }
    except Exception as e:
        return {"error": str(e)}


def extract_archive(path: str, dest: str = None) -> Dict[str, Any]:
    """Extract ZIP, TAR, GZ, or TAR.GZ archives."""
    src_path = Path(path)
    if not src_path.exists():
        return {"error": f"Archive not found: {path}"}
    
    dest_path = Path(dest) if dest else src_path.parent / src_path.stem
    dest_path.mkdir(parents=True, exist_ok=True)
    
    try:
        if src_path.suffix == '.zip':
            import zipfile
            with zipfile.ZipFile(src_path, 'r') as zf:
                zf.extractall(dest_path)
        elif src_path.suffix in ('.tar', '.gz', '.tgz'):
            import tarfile
            mode = 'r:gz' if src_path.suffix in ('.gz', '.tgz') else 'r'
            with tarfile.open(src_path, mode) as tf:
                tf.extractall(dest_path)
        else:
            return {"error": f"Unsupported format: {src_path.suffix}"}
        
        return {"status": "success", "extracted_to": str(dest_path)}
    except Exception as e:
        return {"error": str(e)}


def git_clone(repo_url: str, dest: str, branch: str = None, depth: int = None) -> Dict[str, Any]:
    """Clone a git repository."""
    if not HAVE_GIT:
        return {"error": "GitPython not installed."}
    
    dest_path = Path(dest)
    try:
        import git
        repo = git.Repo.clone_from(repo_url, str(dest_path), branch=branch, depth=depth)
        return {
            "status": "success",
            "path": str(dest_path),
            "branch": branch or repo.active_branch.name,
            "commit": str(repo.head.commit.hexsha)[:8]
        }
    except Exception as e:
        return {"error": str(e)}


def git_pull(path: str) -> Dict[str, Any]:
    """Pull latest changes in a git repository."""
    if not HAVE_GIT:
        return {"error": "GitPython not installed"}
    
    try:
        import git
        repo = git.Repo(path)
        origin = repo.remotes.origin
        pull_info = origin.pull()
        return {
            "status": "success",
            "path": path,
            "after_commit": str(pull_info[0].commit.hexsha)[:8] if pull_info else "up-to-date"
        }
    except Exception as e:
        return {"error": str(e)}


def http_get(url: str, headers: Dict[str, str] = None, timeout: int = 30) -> Dict[str, Any]:
    """Make HTTP GET request."""
    req_headers = headers or {'User-Agent': 'ZeroBound-Agent/1.0'}
    try:
        import urllib.request
        req = urllib.request.Request(url, headers=req_headers)
        with urllib.request.urlopen(req, timeout=timeout) as response:
            content = response.read().decode('utf-8', errors='replace')
            return {
                "status_code": response.status,
                "content": content[:10000],
                "truncated": len(content) > 10000,
                "headers": dict(response.headers)
            }
    except Exception as e:
        return {"error": str(e)}


def http_post(url: str, data: str = None, json_data: Dict = None, headers: Dict[str, str] = None) -> Dict[str, Any]:
    """Make HTTP POST request."""
    req_headers = headers or {'User-Agent': 'ZeroBound-Agent/1.0', 'Content-Type': 'application/json'}
    if json_data:
        body = json.dumps(json_data).encode('utf-8')
    elif data:
        body = data.encode('utf-8')
    else:
        body = b''
    
    try:
        import urllib.request
        req = urllib.request.Request(url, data=body, headers=req_headers, method='POST')
        with urllib.request.urlopen(req, timeout=30) as response:
            content = response.read().decode('utf-8', errors='replace')
            return {
                "status_code": response.status,
                "content": content[:10000],
                "truncated": len(content) > 10000
            }
    except Exception as e:
        return {"error": str(e)}


class FileWatcher:
    """Simple file/directory watcher."""
    def __init__(self):
        self.active_watchers = {}
    
    def watch(self, path: str, recursive: bool = False) -> Dict[str, Any]:
        watcher_id = f"watcher_{int(time.time())}_{hash(path)}"
        def get_snapshot():
            snap = {}
            path_obj = Path(path)
            if path_obj.is_file():
                snap[path] = path_obj.stat().st_mtime
            else:
                pattern = '**/*' if recursive else '*'
                for p in path_obj.glob(pattern):
                    if p.is_file(): snap[str(p)] = p.stat().st_mtime
            return snap
        
        self.active_watchers[watcher_id] = {'path': path, 'initial': get_snapshot(), 'recursive': recursive}
        return {"watcher_id": watcher_id, "message": f"Watching {path}"}
    
    def check_changes(self, watcher_id: str) -> Dict[str, Any]:
        if watcher_id not in self.active_watchers: return {"error": "Watcher not found"}
        w = self.active_watchers[watcher_id]
        path = Path(w['path'])
        new_snap = {}
        if path.is_file():
            if path.exists(): new_snap[str(path)] = path.stat().st_mtime
        else:
            pattern = '**/*' if w['recursive'] else '*'
            for p in path.glob(pattern):
                if p.is_file(): new_snap[str(p)] = p.stat().st_mtime
        
        old = w['initial']
        changes = []
        for f, mt in new_snap.items():
            if f not in old: changes.append({'file': f, 'event': 'created'})
            elif old[f] != mt: changes.append({'file': f, 'event': 'modified'})
        for f in old:
            if f not in new_snap: changes.append({'file': f, 'event': 'deleted'})
        w['initial'] = new_snap
        return {"watcher_id": watcher_id, "changes": changes}

_file_watcher = FileWatcher()

def watch_directory(path: str, recursive: bool = False) -> Dict[str, Any]:
    """Start watching a directory for changes."""
    return _file_watcher.watch(path, recursive)

def check_file_changes(watcher_id: str) -> Dict[str, Any]:
    """Check for file changes."""
    return _file_watcher.check_changes(watcher_id)

def lock_file(path: str, timeout: int = 30) -> Dict[str, Any]:
    """Acquire an advisory lock on a file."""
    lock_path = Path(f"{path}.lock")
    start = time.time()
    while time.time() - start < timeout:
        try:
            if sys.platform == "win32":
                import msvcrt
                f = open(lock_path, 'w')
                msvcrt.locking(f.fileno(), msvcrt.LK_NBLCK, 1)
                return {"status": "locked", "lock_file": str(lock_path)}
            else:
                import fcntl
                f = open(lock_path, 'w')
                fcntl.flock(f, fcntl.LOCK_EX | fcntl.LOCK_NB)
                return {"status": "locked", "lock_file": str(lock_path)}
        except: time.sleep(0.1)
    return {"error": f"Timeout acquiring lock on {path}"}

def get_process_tree(pid: int = None) -> Dict[str, Any]:
    """Get process tree (requires psutil)."""
    if not HAVE_PSUTIL: return {"error": "psutil not installed."}
    pid = pid or os.getpid()
    def get_children(p):
        return [{'pid': c.pid, 'name': c.name(), 'children': get_children(c)} for c in p.children()]
    try:
        import psutil
        proc = psutil.Process(pid)
        return {'pid': pid, 'name': proc.name(), 'children': get_children(proc)}
    except Exception as e: return {"error": str(e)}

def send_signal(pid: int, signal: str) -> Dict[str, Any]:
    """Send a signal to a process (SIGINT, SIGTERM, SIGKILL)."""
    sig_map = {'SIGINT': 2, 'SIGTERM': 15, 'SIGKILL': 9}
    if signal not in sig_map: return {"error": f"Unknown signal: {signal}"}
    try:
        os.kill(pid, sig_map[signal])
        return {"status": "sent", "signal": signal, "pid": pid}
    except Exception as e: return {"error": str(e)}


# =============================================================================
# EXTENDED TOOLS (PRIORITIES 4-7)
# =============================================================================

def pip_install(packages: List[str], upgrade: bool = False) -> Dict[str, Any]:
    """Install Python packages."""
    cmd = [sys.executable, '-m', 'pip', 'install']
    if upgrade: cmd.append('--upgrade')
    cmd.extend(packages)
    try:
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
        return {"status": "success" if result.returncode == 0 else "failed", "stdout": result.stdout[-1000:]}
    except Exception as e: return {"error": str(e)}

def pip_list() -> Dict[str, Any]:
    """List installed packages."""
    try:
        result = subprocess.run([sys.executable, '-m', 'pip', 'list', '--format=json'], capture_output=True, text=True)
        return {"packages": json.loads(result.stdout)}
    except Exception as e: return {"error": str(e)}

def create_virtual_env(path: str) -> Dict[str, Any]:
    """Create a Python virtual environment."""
    try:
        subprocess.run([sys.executable, '-m', 'venv', path], check=True)
        return {"status": "success", "path": path}
    except Exception as e: return {"error": str(e)}

def find_symbol_definition(symbol: str, path: str = ".") -> Dict[str, Any]:
    """Find where a symbol is defined using regex."""
    root = _get_operation_path(path)
    results = []
    patterns = [rf'def\s+{symbol}\s*\(', rf'class\s+{symbol}\s*[:\(]', rf'{symbol}\s*=']
    for py_file in Path(root).rglob('*.py'):
        try:
            with open(py_file, 'r', encoding='utf-8', errors='ignore') as f:
                for i, line in enumerate(f, 1):
                    if any(re.search(p, line) for p in patterns):
                        results.append({'file': str(py_file), 'line': i, 'content': line.strip()})
        except: continue
    return {'symbol': symbol, 'definitions': results[:50]}

def find_all_references(symbol: str, path: str = ".") -> Dict[str, Any]:
    """Find all references to a symbol."""
    root = _get_operation_path(path)
    results = []
    pattern = rf'\b{symbol}\b'
    for py_file in Path(root).rglob('*.py'):
        try:
            with open(py_file, 'r', encoding='utf-8', errors='ignore') as f:
                for i, line in enumerate(f, 1):
                    if re.search(pattern, line):
                        results.append({'file': str(py_file), 'line': i, 'content': line.strip()})
        except: continue
    return {'symbol': symbol, 'references': results[:100]}

def get_imports(file_path: str) -> Dict[str, Any]:
    """Parse Python imports from a file."""
    path = _get_operation_path(file_path)
    imports = []
    try:
        with open(path, 'r', encoding='utf-8') as f:
            for line in f:
                if line.startswith(('import ', 'from ')): imports.append(line.strip())
        return {"file": file_path, "imports": imports}
    except Exception as e: return {"error": str(e)}

def run_pytest_coverage(path: str = ".") -> Dict[str, Any]:
    """Run pytest with coverage."""
    cmd = [sys.executable, '-m', 'pytest', path, '--cov=' + path]
    try:
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
        return {"passed": result.returncode == 0, "stdout": result.stdout[-2000:]}
    except Exception as e: return {"error": str(e)}

def run_linter(path: str) -> Dict[str, Any]:
    """Run ruff linter on the codebase."""
    try:
        result = subprocess.run(['ruff', 'check', path], capture_output=True, text=True)
        return {"output": result.stdout[-2000:]}
    except: return {"error": "ruff not installed."}

def read_data(file_path: str) -> Dict[str, Any]:
    """Read data file into pandas summary."""
    if not HAVE_PANDAS: return {"error": "pandas not installed."}
    path = _get_operation_path(file_path)
    try:
        df = pd.read_csv(path) if path.endswith('.csv') else pd.read_excel(path)
        return {"shape": list(df.shape), "columns": list(df.columns), "head": df.head(5).to_dict()}
    except Exception as e: return {"error": str(e)}

def plot_save(data: Dict, plot_type: str = "line", save_path: str = "plot.png") -> Dict[str, Any]:
    """Create and save a plot."""
    if not HAVE_MATPLOTLIB: return {"error": "matplotlib not installed."}
    try:
        df = pd.DataFrame(data)
        plt.figure(figsize=(10, 6))
        if plot_type == 'line': df.plot()
        elif plot_type == 'bar': df.plot(kind='bar')
        plt.savefig(save_path)
        plt.close()
        return {"status": "success", "saved_to": save_path}
    except Exception as e: return {"error": str(e)}

def run_notebook(path: str, output: str = None) -> Dict[str, Any]:
    """Execute a Jupyter notebook."""
    try:
        import nbformat
        from nbconvert.preprocessors import ExecutePreprocessor
        with open(path) as f:
            nb = nbformat.read(f, as_version=4)
        ep = ExecutePreprocessor(timeout=600, kernel_name='python3')
        ep.preprocess(nb, {'metadata': {'path': os.path.dirname(path)}})
        if output:
            with open(output, 'w') as f: nbformat.write(nb, f)
        return {"status": "success"}
    except Exception as e: return {"error": str(e)}

def create_requirements(path: str = ".", output: str = "requirements.txt") -> Dict[str, Any]:
    """Generate requirements.txt from imports."""
    imports = set()
    for py_file in Path(path).rglob('*.py'):
        try:
            with open(py_file, 'r') as f:
                for line in f:
                    if line.startswith('import '): imports.add(line.split()[1].split('.')[0])
                    elif line.startswith('from '): imports.add(line.split()[1].split('.')[0])
        except: continue
    with open(output, 'w') as f:
        for imp in sorted(imports): f.write(f"{imp}\n")
    return {"status": "success", "file": output}


# ---------------------------------------------------------------------------
# Marimo notebook tools
# ---------------------------------------------------------------------------

def marimo_get_cell_map(notebook_path: str) -> Dict[str, Any]:
    """
    Parse a Marimo .py notebook and return a lightweight map of all cells.
    Returns cell index, first line preview, and approximate line range.
    This is the primary discovery tool — always call this before editing.
    """
    path = _get_operation_path(notebook_path)
    if not os.path.exists(path):
        return {"error": f"Notebook not found: {notebook_path}"}
    try:
        with open(path, "r", encoding="utf-8") as f:
            source = f.read()
        # Marimo cells are decorated with @app.cell or @app.cell(hide_code=True)
        import ast
        tree = ast.parse(source)
        cells = []
        lines = source.splitlines()
        for node in ast.walk(tree):
            if isinstance(node, ast.FunctionDef):
                for decorator in node.decorator_list:
                    dec_src = ast.unparse(decorator) if hasattr(ast, 'unparse') else ""
                    if 'cell' in dec_src or 'app.cell' in dec_src:
                        body_lines = lines[node.body[0].lineno - 1 : node.end_lineno]
                        preview = body_lines[0].strip()[:80] if body_lines else ""
                        cells.append({
                            "index": len(cells),
                            "name": node.name,
                            "line_start": node.lineno,
                            "line_end": node.end_lineno,
                            "preview": preview,
                            "decorator": dec_src
                        })
        return {
            "notebook": notebook_path,
            "cell_count": len(cells),
            "cells": cells
        }
    except Exception as e:
        return {"error": str(e)}


def marimo_get_runtime_data(marimo_server_url: str = "http://localhost:2718") -> Dict[str, Any]:
    """
    Query a running Marimo server for live runtime state:
    variable values, cell statuses (idle/running/error), and stdout.
    Marimo must be running with: marimo edit notebook.py
    """
    try:
        import urllib.request
        import json as json_mod
        req = urllib.request.Request(
            f"{marimo_server_url.rstrip('/')}/api/kernel/status",
            headers={"Accept": "application/json"}
        )
        with urllib.request.urlopen(req, timeout=5) as resp:
            data = json_mod.loads(resp.read().decode())
        return {"status": "ok", "runtime": data}
    except Exception as e:
        return {
            "status": "unavailable",
            "hint": "Start Marimo with: marimo edit notebook.py. Then retry.",
            "error": str(e)
        }


def marimo_update_cell(notebook_path: str, cell_name: str, new_code: str) -> Dict[str, Any]:
    """
    Replace the body of a named Marimo cell in the notebook .py file.
    The cell is identified by its function name (from marimo_get_cell_map).
    If Marimo is running with --watch, changes apply instantly in the browser.
    Uses the same safe edit_file fuzzy-match logic internally.
    """
    path = _get_operation_path(notebook_path)
    if not os.path.exists(path):
        return {"error": f"Notebook not found: {notebook_path}"}
    try:
        with open(path, "r", encoding="utf-8") as f:
            source = f.read()
        import ast
        tree = ast.parse(source)
        lines = source.splitlines(keepends=True)
        target_node = None
        for node in ast.walk(tree):
            if isinstance(node, ast.FunctionDef) and node.name == cell_name:
                for dec in node.decorator_list:
                    dec_src = ast.unparse(dec) if hasattr(ast, 'unparse') else ""
                    if 'cell' in dec_src:
                        target_node = node
                        break
            if target_node:
                break
        if not target_node:
            return {"error": f"Cell '{cell_name}' not found in {notebook_path}"}
        # Preserve function signature and decorator; replace only the body
        body_start = target_node.body[0].lineno - 1   # 0-indexed
        body_end = target_node.end_lineno              # inclusive
        # Detect indentation from first body line
        indent = ""
        first_body_line = lines[body_start] if body_start < len(lines) else ""
        for ch in first_body_line:
            if ch in (" ", "\t"): indent += ch
            else: break
        # Indent the new code body
        new_body_lines = [(indent + l if l.strip() else l) for l in new_code.splitlines()]
        new_lines = lines[:body_start] + [l + "\n" for l in new_body_lines] + lines[body_end:]
        new_source = "".join(new_lines)
        # Validate syntax before writing
        syntax_err = check_python_syntax(new_source)
        if syntax_err:
            return {"error": "Syntax error in new_code", "details": syntax_err}
        with open(path, "w", encoding="utf-8") as f:
            f.write(new_source)
        _file_cache.pop(path, None)
        return {"status": "success", "cell": cell_name, "notebook": notebook_path}
    except Exception as e:
        return {"error": str(e)}


def marimo_add_cell(notebook_path: str, code: str, after_cell: str = None, hide_code: bool = False) -> Dict[str, Any]:
    """
    Inject a new Marimo cell into the notebook .py file.
    Inserts after the specified cell (by name), or appends at the end.
    The cell is given an auto-generated unique name.
    If Marimo is running with --watch, the new cell appears live in the browser.
    """
    path = _get_operation_path(notebook_path)
    if not os.path.exists(path):
        return {"error": f"Notebook not found: {notebook_path}"}
    try:
        with open(path, "r", encoding="utf-8") as f:
            source = f.read()
        import ast, uuid as _uuid
        cell_name = "_cell_" + _uuid.uuid4().hex[:8]
        decorator = "@app.cell(hide_code=True)" if hide_code else "@app.cell"
        indent = "    "  # Marimo uses 4-space indent inside cells
        indented_code = "\n".join(indent + l if l.strip() else "" for l in code.splitlines())
        new_cell = f"\n\n{decorator}\ndef {cell_name}():\n{indented_code}\n"
        if after_cell:
            tree = ast.parse(source)
            lines = source.splitlines(keepends=True)
            insert_line = None
            for node in ast.walk(tree):
                if isinstance(node, ast.FunctionDef) and node.name == after_cell:
                    for dec in node.decorator_list:
                        dec_src = ast.unparse(dec) if hasattr(ast, 'unparse') else ""
                        if 'cell' in dec_src:
                            insert_line = node.end_lineno  # 1-indexed
                            break
            if insert_line is None:
                return {"error": f"Cell '{after_cell}' not found; cannot insert after it."}
            new_lines = lines[:insert_line] + [new_cell] + lines[insert_line:]
            new_source = "".join(new_lines)
        else:
            # Append before the final `app.run()` call if present, else at end
            if "app.run()" in source:
                new_source = source.replace("app.run()", new_cell + "\napp.run()")
            else:
                new_source = source + new_cell
        syntax_err = check_python_syntax(new_source)
        if syntax_err:
            return {"error": "Syntax error in new cell code", "details": syntax_err}
        with open(path, "w", encoding="utf-8") as f:
            f.write(new_source)
        _file_cache.pop(path, None)
        return {"status": "success", "new_cell_name": cell_name, "notebook": notebook_path}
    except Exception as e:
        return {"error": str(e)}


def marimo_check_notebook(notebook_path: str) -> Dict[str, Any]:
    """
    Validate a Marimo notebook for syntax and circular dependency errors
    without running it. Equivalent to: marimo check notebook.py
    Returns any errors found so the agent can self-correct.
    """
    path = _get_operation_path(notebook_path)
    if not os.path.exists(path):
        return {"error": f"Notebook not found: {notebook_path}"}
    # 1. Python syntax check
    with open(path, "r", encoding="utf-8") as f:
        source = f.read()
    syntax_err = check_python_syntax(source)
    if syntax_err:
        return {"valid": False, "error_type": "syntax", "details": syntax_err}
    # 2. Try running 'marimo check' via CLI if available
    try:
        result = subprocess.run(
            ["marimo", "check", path],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode != 0:
            return {
                "valid": False,
                "error_type": "marimo_check",
                "stdout": result.stdout,
                "stderr": result.stderr
            }
        return {"valid": True, "message": result.stdout or "No errors found."}
    except FileNotFoundError:
        return {
            "valid": None,
            "warning": "'marimo' CLI not found in PATH. Only Python syntax was checked.",
            "python_syntax": "ok"
        }
    except Exception as e:
        return {"error": str(e)}


def marimo_validate(notebook_path: str) -> Dict[str, Any]:
    """
    Extended validation: checks for duplicate global variable definitions across cells,
    which would cause Marimo's DAG to break. Returns the list of conflicts found.
    Always run this after adding or updating cells that define new variables.
    """
    path = _get_operation_path(notebook_path)
    if not os.path.exists(path):
        return {"error": f"Notebook not found: {notebook_path}"}
    try:
        import ast
        with open(path, "r", encoding="utf-8") as f:
            source = f.read()
        tree = ast.parse(source)
        # Track variable definitions per cell
        cell_defs: Dict[str, list] = {}  # var_name -> [cell_names]
        for node in ast.walk(tree):
            if isinstance(node, ast.FunctionDef):
                for dec in node.decorator_list:
                    dec_src = ast.unparse(dec) if hasattr(ast, 'unparse') else ""
                    if 'cell' in dec_src:
                        # Find all top-level assignments in the cell body
                        for stmt in node.body:
                            if isinstance(stmt, ast.Assign):
                                for target in stmt.targets:
                                    if isinstance(target, ast.Name):
                                        cell_defs.setdefault(target.id, []).append(node.name)
                            elif isinstance(stmt, (ast.AnnAssign,)):
                                if isinstance(stmt.target, ast.Name):
                                    cell_defs.setdefault(stmt.target.id, []).append(node.name)
        # Find conflicts (same var defined in >1 cell)
        conflicts = {var: cells for var, cells in cell_defs.items() if len(cells) > 1}
        if conflicts:
            return {
                "valid": False,
                "conflicts": conflicts,
                "hint": "Each global variable must be defined in exactly ONE cell. Merge or rename duplicates."
            }
        return {"valid": True, "message": f"No duplicate variable definitions found. {len(cell_defs)} globals checked."}
    except Exception as e:
        return {"error": str(e)}


# ---------------------------------------------------------------------------
# Tool dispatcher
# ---------------------------------------------------------------------------
TOOL_MAP = {
    "set_workspace": (set_workspace, False),
    "create_folder": (create_folder, False),
    "run_command": (run_command_async, True),      # async
    "run_shell_command": (run_command_async, True),
    "start_background_command": (start_background_command, True),
    "read_process_output": (read_process_output, False),
    "kill_process": (kill_process, False),
    "list_running_processes": (list_running_processes, False),
    "write_file": (write_file, False),
    "edit_file": (edit_file, False),
    "append_file": (append_file, False),
    "read_file": (read_file, False),
    "read_files": (read_files, False),
    "grep_search": (grep_search, False),
    "find_files": (find_files, False),
    "list_files": (list_files, False),
    "get_file_tree": (get_file_tree, False),
    "reveal_in_os": (reveal_in_os, False),
    "delete_file": (delete_file, False),
    "move_file": (move_file, False),
    "copy_file": (copy_file, False),
    "get_file_info": (get_file_info, False),
    "resolve_path": (resolve_path, False),
    "diagnose_path": (diagnose_path, False),
    "verify_junction": (verify_junction, False),
    "resolve_junction": (resolve_junction, False),
    "search_web": (search_web, True),
    "read_url": (read_url, True),
    "store_memory": (store_memory, False),
    "recall_memory": (recall_memory, False),
    "git_commit": (git_commit, False),
    "git_diff": (git_diff, False),
    "run_tests": (run_tests, False),
    "get_definition": (get_definition, False),
    
    # Extended Tools (Priority 1-7)
    "get_env_var": (get_env_var, False),
    "set_env_var": (set_env_var, False),
    "get_system_info": (get_system_info, False),
    "is_admin": (is_admin, False),
    "download_file": (download_file, False),
    "extract_archive": (extract_archive, False),
    "git_clone": (git_clone, False),
    "git_pull": (git_pull, False),
    "http_get": (http_get, False),
    "http_post": (http_post, False),
    "watch_directory": (watch_directory, False),
    "check_file_changes": (check_file_changes, False),
    "lock_file": (lock_file, False),
    "get_process_tree": (get_process_tree, False),
    "send_signal": (send_signal, False),
    "pip_install": (pip_install, False),
    "pip_list": (pip_list, False),
    "create_virtual_env": (create_virtual_env, False),
    "find_symbol_definition": (find_symbol_definition, False),
    "find_all_references": (find_all_references, False),
    "get_imports": (get_imports, False),
    "run_pytest_coverage": (run_pytest_coverage, False),
    "run_linter": (run_linter, False),
    "read_data": (read_data, False),
    "plot_save": (plot_save, False),
    "run_notebook": (run_notebook, False),
    "create_requirements": (create_requirements, False),
    # Marimo notebook tools
    "marimo_get_cell_map": (marimo_get_cell_map, False),
    "marimo_get_runtime_data": (marimo_get_runtime_data, False),
    "marimo_update_cell": (marimo_update_cell, False),
    "marimo_add_cell": (marimo_add_cell, False),
    "marimo_check_notebook": (marimo_check_notebook, False),
    "marimo_validate": (marimo_validate, False),
}


async def handle_tool_call(tool_name: str, args: Dict[str, Any], callback: Optional[Callable] = None) -> Any:
    """Dispatch tool call, automatically handling async tools."""
    if tool_name.startswith("browser_"):
        bm = _get_browser_manager()
        cmd = tool_name.split("_", 1)[1]
        if cmd == "goto":
            return await bm.goto(args["url"])
        elif cmd == "click":
            return await bm.click(args["selector"])
        elif cmd == "type":
            return await bm.type(args["selector"], args["text"])
        elif cmd == "scroll":
            return await bm.scroll(args.get("direction", "down"), args.get("amount", 500))
        elif cmd == "screenshot":
            return await bm.screenshot()
        elif cmd == "close":
            return await bm.close()
        else:
            return {"error": f"Unknown browser command: {cmd}"}

    if tool_name not in TOOL_MAP:
        return {"error": f"Tool '{tool_name}' not found"}
    func, is_async = TOOL_MAP[tool_name]
    if is_async:
        return await func(**{k: v for k, v in args.items() if v is not None}, callback=callback)
    else:
        return func(**{k: v for k, v in args.items() if v is not None})


# ---------------------------------------------------------------------------
# Tool schema definitions (for LLM prompt)
# ---------------------------------------------------------------------------
TOOLS = [
    {
        "name": "read_file",
        "description": "Read content from a file (or image as base64).",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Relative or absolute path to the file."},
                "start_line": {"type": "integer", "description": "Optional: Start line (1-indexed)."},
                "end_line": {"type": "integer", "description": "Optional: End line (inclusive)."}
            },
            "required": ["path"]
        }
    },
    {
        "name": "read_files",
        "description": "Read multiple small/medium files at once. Best for batch inspection of codebases.",
        "parameters": {
            "type": "object",
            "properties": {
                "paths": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "List of file paths to read."
                }
            },
            "required": ["paths"]
        }
    },
    {
        "name": "write_file",
        "description": "Write or overwrite a file with new content.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Relative or absolute path."},
                "content": {"type": "string", "description": "Full file content."},
                "content_base64": {"type": "string", "description": "Base64-encoded file content (preserves all whitespace)."},
                "lines": {"type": "array", "items": {"type": "string"}, "description": "Alternative to content: Provide file content as an array of lines."},
                "lines_base64": {"type": "array", "items": {"type": "string"}, "description": "Alternative to content: array of base64 lines."}
            },
            "required": ["path"]
        }
    },
    {
        "name": "edit_file",
        "description": "Surgically edit a file by replacing a unique block of text.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Path to the file."},
                "target": {"type": "string", "description": "The exact string to be replaced."},
                "replacement": {"type": "string", "description": "The new content."},
                "target_base64": {"type": "string", "description": "Base64-encoded target."},
                "replacement_base64": {"type": "string", "description": "Base64-encoded replacement."},
                "target_lines": {"type": "array", "items": {"type": "string"}, "description": "Alternative to target: array of lines."},
                "replacement_lines": {"type": "array", "items": {"type": "string"}, "description": "Alternative to replacement: array of lines."},
                "target_lines_base64": {"type": "array", "items": {"type": "string"}, "description": "Alternative to target: array of base64 lines."},
                "replacement_lines_base64": {"type": "array", "items": {"type": "string"}, "description": "Alternative to replacement: array of base64 lines."}
            },
            "required": ["path"]
        }
    },
    {
        "name": "append_file",
        "description": "Append content to the end of a file.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "content": {"type": "string"},
                "content_base64": {"type": "string", "description": "Base64-encoded content."},
                "lines": {"type": "array", "items": {"type": "string"}, "description": "Alternative to content: array of lines."},
                "lines_base64": {"type": "array", "items": {"type": "string"}, "description": "Alternative to content: array of base64 lines."}
            },
            "required": ["path"]
        }
    },
    {
        "name": "create_folder",
        "description": "Create a new directory (recursively).",
        "parameters": {
            "type": "object",
            "properties": {"path": {"type": "string"}},
            "required": ["path"]
        }
    },
    {
        "name": "delete_file",
        "description": "Delete a file or folder.",
        "parameters": {
            "type": "object",
            "properties": {"path": {"type": "string"}},
            "required": ["path"]
        }
    },
    {
        "name": "move_file",
        "description": "Move or rename a file or folder.",
        "parameters": {
            "type": "object",
            "properties": {
                "source": {"type": "string", "description": "Current path."},
                "destination": {"type": "string", "description": "New path."}
            },
            "required": ["source", "destination"]
        }
    },
    {
        "name": "copy_file",
        "description": "Copy a file or folder.",
        "parameters": {
            "type": "object",
            "properties": {
                "source": {"type": "string", "description": "Source path."},
                "destination": {"type": "string", "description": "Destination path."}
            },
            "required": ["source", "destination"]
        }
    },
    {
        "name": "diagnose_path",
        "description": "Provide deep diagnostic info about path visibility for debugging.",
        "parameters": {
            "type": "object",
            "properties": {"path": {"type": "string"}},
            "required": ["path"]
        }
    },
    {
        "name": "list_files",
        "description": "List files in a directory.",
        "parameters": {
            "type": "object",
            "properties": {"path": {"type": "string", "default": "."}}
        }
    },
    {
        "name": "get_file_tree",
        "description": "Get a recursive tree view of the workspace.",
        "parameters": {
            "type": "object",
            "properties": {"startpath": {"type": "string", "description": "Optional starting directory."}}
        }
    },
    {
        "name": "get_file_info",
        "description": "Get metadata about a file or folder.",
        "parameters": {
            "type": "object",
            "properties": {"path": {"type": "string"}},
            "required": ["path"]
        }
    },
    {
        "name": "resolve_path",
        "description": "Resolve a logical path to its actual physical filesystem path, accounting for any internal mappings.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "The path to resolve (e.g., '.', './folder', 'D:\\Downloads\\...')."}
            },
            "required": ["path"]
        }
    },
    {
        "name": "grep_search",
        "description": "Search for a pattern in files (regex).",
        "parameters": {
            "type": "object",
            "properties": {
                "pattern": {"type": "string"},
                "path": {"type": "string", "default": "."},
                "case_insensitive": {"type": "boolean", "default": True}
            },
            "required": ["pattern"]
        }
    },
    {
        "name": "find_files",
        "description": "Find files by name pattern (wildcard supported).",
        "parameters": {
            "type": "object",
            "properties": {"pattern": {"type": "string"}},
            "required": ["pattern"]
        }
    },
    {
        "name": "reveal_in_os",
        "description": "Open the file/folder in system file explorer.",
        "parameters": {
            "type": "object",
            "properties": {"path": {"type": "string"}},
            "required": ["path"]
        }
    },
    {
        "name": "run_command",
        "description": "Execute a shell command (async, streams to terminal, auto-timeout 15s).",
        "parameters": {
            "type": "object",
            "properties": {"command": {"type": "string"}},
            "required": ["command"]
        }
    },
    {
        "name": "start_background_command",
        "description": "Start a long-running background process (no timeout).",
        "parameters": {
            "type": "object",
            "properties": {"command": {"type": "string"}},
            "required": ["command"]
        }
    },
    {
        "name": "read_process_output",
        "description": "Read new output from a background process.",
        "parameters": {
            "type": "object",
            "properties": {"process_id": {"type": "string"}},
            "required": ["process_id"]
        }
    },
    {
        "name": "kill_process",
        "description": "Terminate a background process.",
        "parameters": {
            "type": "object",
            "properties": {"process_id": {"type": "string"}},
            "required": ["process_id"]
        }
    },
    {
        "name": "list_running_processes",
        "description": "List all currently running background processes.",
        "parameters": {
            "type": "object",
            "properties": {}
        }
    },
    {
        "name": "search_web",
        "description": "Search the web using DuckDuckGo.",
        "parameters": {
            "type": "object",
            "properties": {"query": {"type": "string"}},
            "required": ["query"]
        }
    },
    {
        "name": "read_url",
        "description": "Fetch and clean text from a URL.",
        "parameters": {
            "type": "object",
            "properties": {"url": {"type": "string"}},
            "required": ["url"]
        }
    },
    {
        "name": "browser_goto",
        "description": "Navigate the automated browser to a URL.",
        "parameters": {
            "type": "object",
            "properties": {"url": {"type": "string"}},
            "required": ["url"]
        }
    },
    {
        "name": "browser_click",
        "description": "Click an element in the browser.",
        "parameters": {
            "type": "object",
            "properties": {"selector": {"type": "string"}},
            "required": ["selector"]
        }
    },
    {
        "name": "browser_type",
        "description": "Type text into a browser input.",
        "parameters": {
            "type": "object",
            "properties": {
                "selector": {"type": "string"},
                "text": {"type": "string"}
            },
            "required": ["selector", "text"]
        }
    },
    {
        "name": "browser_scroll",
        "description": "Scroll the browser page.",
        "parameters": {
            "type": "object",
            "properties": {
                "direction": {"type": "string", "enum": ["up", "down"], "default": "down"},
                "amount": {"type": "integer", "default": 500}
            }
        }
    },
    {
        "name": "browser_screenshot",
        "description": "Take a screenshot of the current browser page.",
        "parameters": {"type": "object", "properties": {}}
    },
    {
        "name": "browser_close",
        "description": "Close the browser instance.",
        "parameters": {"type": "object", "properties": {}}
    },
    {
        "name": "store_memory",
        "description": "Store a key-value pair in long-term memory.",
        "parameters": {
            "type": "object",
            "properties": {
                "key": {"type": "string"},
                "value": {"type": "string"}
            },
            "required": ["key", "value"]
        }
    },
    {
        "name": "recall_memory",
        "description": "Search for memories by key or value.",
        "parameters": {
            "type": "object",
            "properties": {"query": {"type": "string"}},
            "required": ["query"]
        }
    },
    {
        "name": "git_commit",
        "description": "Commit current changes to git.",
        "parameters": {
            "type": "object",
            "properties": {
                "message": {"type": "string"},
                "files": {"type": "array", "items": {"type": "string"}, "description": "Optional: specific files to commit"}
            },
            "required": ["message"]
        }
    },
    {
        "name": "git_diff",
        "description": "Show unstaged changes (or specific file diff).",
        "parameters": {
            "type": "object",
            "properties": {"path": {"type": "string", "description": "Optional: specific file path"}}
        }
    },
    {
        "name": "run_tests",
        "description": "Run pytest tests.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "default": "."},
                "pattern": {"type": "string", "description": "Optional: test name pattern (-k)"}
            }
        }
    },
    {
        "name": "get_definition",
        "description": "Get definition of symbol at position using Jedi.",
        "parameters": {
            "type": "object",
            "properties": {
                "file": {"type": "string"},
                "line": {"type": "integer"},
                "column": {"type": "integer"}
            },
            "required": ["file", "line", "column"]
        }
    },
    {
        "name": "set_workspace",
        "description": "Change the active workspace directory.",
        "parameters": {
            "type": "object",
            "properties": {"path": {"type": "string"}},
            "required": ["path"]
        }
    },
    {
        "name": "get_knowledge_stats",
        "description": "Get statistics from the pattern learning system.",
        "parameters": {"type": "object", "properties": {}}
    },
    {
        "name": "get_env_var",
        "description": "Get environment variable value.",
        "parameters": {
            "type": "object",
            "properties": {"name": {"type": "string"}},
            "required": ["name"]
        }
    },
    {
        "name": "set_env_var",
        "description": "Set environment variable (permanent requires admin).",
        "parameters": {
            "type": "object",
            "properties": {
                "name": {"type": "string"},
                "value": {"type": "string"},
                "permanent": {"type": "boolean", "default": False}
            },
            "required": ["name", "value"]
        }
    },
    {
        "name": "get_system_info",
        "description": "Get comprehensive system information (OS, CPU, RAM, Disk).",
        "parameters": {"type": "object", "properties": {}}
    },
    {
        "name": "download_file",
        "description": "Download a file from URL with resume support.",
        "parameters": {
            "type": "object",
            "properties": {
                "url": {"type": "string"},
                "dest": {"type": "string"},
                "resume": {"type": "boolean", "default": False}
            },
            "required": ["url", "dest"]
        }
    },
    {
        "name": "extract_archive",
        "description": "Extract ZIP or TAR archives.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "dest": {"type": "string", "description": "Optional destination folder"}
            },
            "required": ["path"]
        }
    },
    {
        "name": "git_clone",
        "description": "Clone a git repository.",
        "parameters": {
            "type": "object",
            "properties": {
                "repo_url": {"type": "string"},
                "dest": {"type": "string"},
                "branch": {"type": "string", "default": None},
                "depth": {"type": "integer", "default": None}
            },
            "required": ["repo_url", "dest"]
        }
    },
    {
        "name": "git_pull",
        "description": "Pull latest changes in a git repository.",
        "parameters": {
            "type": "object",
            "properties": {"path": {"type": "string"}},
            "required": ["path"]
        }
    },
    {
        "name": "http_get",
        "description": "Make HTTP GET request.",
        "parameters": {
            "type": "object",
            "properties": {
                "url": {"type": "string"},
                "headers": {"type": "object", "default": None}
            },
            "required": ["url"]
        }
    },
    {
        "name": "http_post",
        "description": "Make HTTP POST request.",
        "parameters": {
            "type": "object",
            "properties": {
                "url": {"type": "string"},
                "json_data": {"type": "object", "default": None},
                "headers": {"type": "object", "default": None}
            },
            "required": ["url"]
        }
    },
    {
        "name": "watch_directory",
        "description": "Start watching a directory for changes.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "recursive": {"type": "boolean", "default": False}
            },
            "required": ["path"]
        }
    },
    {
        "name": "check_file_changes",
        "description": "Check for changes in a watched directory.",
        "parameters": {
            "type": "object",
            "properties": {"watcher_id": {"type": "string"}},
            "required": ["watcher_id"]
        }
    },
    {
        "name": "pip_install",
        "description": "Install Python packages.",
        "parameters": {
            "type": "object",
            "properties": {
                "packages": {"type": "array", "items": {"type": "string"}},
                "upgrade": {"type": "boolean", "default": False}
            },
            "required": ["packages"]
        }
    },
    {
        "name": "pip_list",
        "description": "List installed Python packages.",
        "parameters": {"type": "object", "properties": {}}
    },
    {
        "name": "find_symbol_definition",
        "description": "Find where a symbol is defined across files.",
        "parameters": {
            "type": "object",
            "properties": {
                "symbol": {"type": "string"},
                "path": {"type": "string", "default": "."}
            },
            "required": ["symbol"]
        }
    },
    {
        "name": "find_all_references",
        "description": "Find all references to a symbol across files.",
        "parameters": {
            "type": "object",
            "properties": {
                "symbol": {"type": "string"},
                "path": {"type": "string", "default": "."}
            },
            "required": ["symbol"]
        }
    },
    {
        "name": "get_imports",
        "description": "Extract imports from a Python file.",
        "parameters": {
            "type": "object",
            "properties": {"file_path": {"type": "string"}},
            "required": ["file_path"]
        }
    },
    {
        "name": "run_pytest_coverage",
        "description": "Run pytest with coverage reporting.",
        "parameters": {
            "type": "object",
            "properties": {"path": {"type": "string", "default": "."}},
            "required": ["path"]
        }
    },
    {
        "name": "run_linter",
        "description": "Run ruff linter on a path.",
        "parameters": {
            "type": "object",
            "properties": {"path": {"type": "string"}},
            "required": ["path"]
        }
    },
    {
        "name": "read_data",
        "description": "Read data (CSV/Excel) and show summary statistics.",
        "parameters": {
            "type": "object",
            "properties": {"file_path": {"type": "string"}},
            "required": ["file_path"]
        }
    },
    {
        "name": "plot_save",
        "description": "Create and save a plot from data.",
        "parameters": {
            "type": "object",
            "properties": {
                "data": {"type": "object", "description": "Dictionary of lists"},
                "plot_type": {"type": "string", "enum": ["line", "bar"], "default": "line"},
                "save_path": {"type": "string", "default": "plot.png"}
            },
            "required": ["data"]
        }
    },
    {
        "name": "run_notebook",
        "description": "Execute a Jupyter notebook.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "output": {"type": "string", "description": "Optional output path"}
            },
            "required": ["path"]
        }
    },
    {
        "name": "create_requirements",
        "description": "Generate requirements.txt from project imports.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "default": "."},
                "output": {"type": "string", "default": "requirements.txt"}
            }
        }
    },
    {
        "name": "marimo_get_cell_map",
        "description": "Parse a Marimo .py notebook and return a lightweight map of all cells (name, line range, preview). ALWAYS call this first before editing any Marimo notebook to discover existing definitions and avoid duplicate variable conflicts.",
        "parameters": {
            "type": "object",
            "properties": {
                "notebook_path": {"type": "string", "description": "Path to the Marimo .py notebook file."}
            },
            "required": ["notebook_path"]
        }
    },
    {
        "name": "marimo_get_runtime_data",
        "description": "Query a running Marimo server for live cell statuses, variable state, and stdout. Requires Marimo to be running with: marimo edit notebook.py",
        "parameters": {
            "type": "object",
            "properties": {
                "marimo_server_url": {"type": "string", "description": "URL of the Marimo server. Default: http://localhost:2718"}
            }
        }
    },
    {
        "name": "marimo_update_cell",
        "description": "Replace the body of a named Marimo cell with new code. The cell name comes from marimo_get_cell_map. If Marimo is running with --watch, changes are reflected live in the browser. CRITICAL: Never define a variable in this cell that is already defined globally in another cell.",
        "parameters": {
            "type": "object",
            "properties": {
                "notebook_path": {"type": "string", "description": "Path to the Marimo .py notebook."},
                "cell_name": {"type": "string", "description": "The Python function name of the cell to update (from marimo_get_cell_map)."},
                "new_code": {"type": "string", "description": "The new Python code for the cell body."}
            },
            "required": ["notebook_path", "cell_name", "new_code"]
        }
    },
    {
        "name": "marimo_add_cell",
        "description": "Inject a new cell into a Marimo notebook. Inserts after a named cell or appends to the end. Use marimo_validate after adding to check for duplicate variable conflicts.",
        "parameters": {
            "type": "object",
            "properties": {
                "notebook_path": {"type": "string", "description": "Path to the Marimo .py notebook."},
                "code": {"type": "string", "description": "Python code for the new cell body."},
                "after_cell": {"type": "string", "description": "Optional: Insert after this cell name. If omitted, appends before app.run()."},
                "hide_code": {"type": "boolean", "description": "If true, hides the code in the notebook UI (good for UI-only cells)."}
            },
            "required": ["notebook_path", "code"]
        }
    },
    {
        "name": "marimo_check_notebook",
        "description": "Validate a Marimo notebook for syntax and circular dependency errors (runs 'marimo check'). Returns any errors so you can self-correct. Run this after every set of edits.",
        "parameters": {
            "type": "object",
            "properties": {
                "notebook_path": {"type": "string", "description": "Path to the Marimo .py notebook."}
            },
            "required": ["notebook_path"]
        }
    },
    {
        "name": "marimo_validate",
        "description": "Check a Marimo notebook for duplicate global variable definitions across cells — the most common cause of DAG breakage. Always run after adding or modifying cells that define new variables.",
        "parameters": {
            "type": "object",
            "properties": {
                "notebook_path": {"type": "string", "description": "Path to the Marimo .py notebook."}
            },
            "required": ["notebook_path"]
        }
    }
]

def get_tools_prompt_description() -> str:
    """Generate a human-readable description of all tools for the prompt."""
    lines = ["--- AVAILABLE TOOLS ---"]
    for t in TOOLS:
        name = t.get("name", "unknown")
        desc = t.get("description", "")
        params = t.get("parameters", {}).get("properties", {})
        required = t.get("parameters", {}).get("required", [])
        
        p_list = []
        for param_name, param_info in params.items():
            param_desc = param_info.get("description", "")
            if param_name in required:
                p_list.append(f"{param_name} (required{(': ' + param_desc) if param_desc else ''})")
            else:
                p_list.append(f"{param_name} (optional{(': ' + param_desc) if param_desc else ''})")
        
        if p_list:
            lines.append(f"• {name}({', '.join(p_list)}): {desc}")
        else:
            lines.append(f"• {name}(): {desc}")
    
    return "\n".join(lines)